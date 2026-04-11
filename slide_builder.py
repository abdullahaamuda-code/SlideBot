import os
import uuid
import random
import requests
from io import BytesIO

from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from PIL import Image, ImageDraw, ImageFilter

load_dotenv()

UNSPLASH_KEY = os.getenv("UNSPLASH_ACCESS_KEY")
PIXABAY_KEY  = os.getenv("PIXABAY_API_KEY")

# ─────────────────────────────────────────────────────────────────
#  DESIGN SYSTEM
#  Rule: one dominant dark, one light bg, one sharp accent.
#  Dark slides (title / conclusion) anchor the "sandwich".
#  Content slides breathe on light bg with accent pops.
# ─────────────────────────────────────────────────────────────────
THEMES = {
    "classic": {
        # Midnight Executive
        "dark_bg":      RGBColor(0x1E, 0x27, 0x61),   # deep navy
        "light_bg":     RGBColor(0xF7, 0xF9, 0xFF),   # near-white ice
        "accent":       RGBColor(0x4A, 0x90, 0xD9),   # sky blue
        "accent_soft":  RGBColor(0xCA, 0xDC, 0xFC),   # pale ice blue
        "heading":      RGBColor(0x1E, 0x27, 0x61),   # navy
        "body":         RGBColor(0x2D, 0x2D, 0x2D),   # near-black
        "white":        RGBColor(0xFF, 0xFF, 0xFF),
        "muted":        RGBColor(0x8A, 0x9B, 0xBF),
        "card_bg":      RGBColor(0xE8, 0xEF, 0xFB),
    },
    "dark": {
        # Obsidian & Coral
        "dark_bg":      RGBColor(0x12, 0x12, 0x1E),
        "light_bg":     RGBColor(0x1C, 0x1C, 0x2E),
        "accent":       RGBColor(0xFF, 0x5F, 0x6D),   # coral
        "accent_soft":  RGBColor(0xFF, 0x9A, 0x9E),
        "heading":      RGBColor(0xFF, 0xFF, 0xFF),
        "body":         RGBColor(0xCC, 0xCC, 0xDD),
        "white":        RGBColor(0xFF, 0xFF, 0xFF),
        "muted":        RGBColor(0x66, 0x66, 0x88),
        "card_bg":      RGBColor(0x24, 0x24, 0x3E),
    },
    "corporate": {
        # Ocean Authority
        "dark_bg":      RGBColor(0x06, 0x5A, 0x82),
        "light_bg":     RGBColor(0xF2, 0xF7, 0xFB),
        "accent":       RGBColor(0x02, 0xC3, 0x9A),   # teal-mint
        "accent_soft":  RGBColor(0xC8, 0xF0, 0xE8),
        "heading":      RGBColor(0x06, 0x5A, 0x82),
        "body":         RGBColor(0x1A, 0x1A, 0x2E),
        "white":        RGBColor(0xFF, 0xFF, 0xFF),
        "muted":        RGBColor(0x7A, 0xA8, 0xC4),
        "card_bg":      RGBColor(0xE0, 0xF2, 0xFB),
    },
    "startup": {
        # Coral Energy
        "dark_bg":      RGBColor(0x2F, 0x3C, 0x7E),   # deep indigo
        "light_bg":     RGBColor(0xFF, 0xFD, 0xFA),
        "accent":       RGBColor(0xF9, 0x61, 0x67),   # coral-red
        "accent_soft":  RGBColor(0xFD, 0xD5, 0xD6),
        "heading":      RGBColor(0x2F, 0x3C, 0x7E),
        "body":         RGBColor(0x22, 0x22, 0x33),
        "white":        RGBColor(0xFF, 0xFF, 0xFF),
        "muted":        RGBColor(0x9A, 0x9A, 0xBB),
        "card_bg":      RGBColor(0xFE, 0xED, 0xED),
    },
    "academic": {
        # Forest & Moss
        "dark_bg":      RGBColor(0x2C, 0x5F, 0x2D),   # forest green
        "light_bg":     RGBColor(0xF6, 0xFB, 0xF4),
        "accent":       RGBColor(0x97, 0xBC, 0x62),   # moss
        "accent_soft":  RGBColor(0xD7, 0xEB, 0xBB),
        "heading":      RGBColor(0x2C, 0x5F, 0x2D),
        "body":         RGBColor(0x1E, 0x2E, 0x1E),
        "white":        RGBColor(0xFF, 0xFF, 0xFF),
        "muted":        RGBColor(0x7A, 0xA0, 0x7B),
        "card_bg":      RGBColor(0xE2, 0xF3, 0xD9),
    },
    "minimal": {
        # Charcoal Minimal
        "dark_bg":      RGBColor(0x21, 0x21, 0x21),
        "light_bg":     RGBColor(0xFA, 0xFA, 0xFA),
        "accent":       RGBColor(0x21, 0x21, 0x21),
        "accent_soft":  RGBColor(0xE0, 0xE0, 0xE0),
        "heading":      RGBColor(0x21, 0x21, 0x21),
        "body":         RGBColor(0x33, 0x33, 0x33),
        "white":        RGBColor(0xFF, 0xFF, 0xFF),
        "muted":        RGBColor(0x99, 0x99, 0x99),
        "card_bg":      RGBColor(0xEE, 0xEE, 0xEE),
    },
}

FREE_THEMES    = ["classic", "dark"]
PREMIUM_THEMES = ["corporate", "startup", "academic", "minimal"]

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ─────────────────────────────────────────────────────────────────
#  IMAGE FETCHING
# ─────────────────────────────────────────────────────────────────
def fetch_unsplash_image(keyword: str):
    try:
        if not UNSPLASH_KEY:
            return None
        r = requests.get(
            "https://api.unsplash.com/photos/random",
            params={"query": keyword, "orientation": "landscape",
                    "content_filter": "high", "client_id": UNSPLASH_KEY},
            timeout=12,
        )
        if r.status_code == 200:
            img_url = r.json()["urls"]["regular"]
            img_r = requests.get(img_url, timeout=12)
            if img_r.status_code == 200:
                return BytesIO(img_r.content)
    except Exception as e:
        print(f"Unsplash error: {e}")
    return None


def fetch_pixabay_image(keyword: str):
    try:
        if not PIXABAY_KEY:
            return None
        r = requests.get(
            "https://pixabay.com/api/",
            params={"key": PIXABAY_KEY, "q": keyword, "image_type": "photo",
                    "orientation": "horizontal", "safesearch": "true", "per_page": 5},
            timeout=12,
        )
        if r.status_code == 200:
            hits = r.json().get("hits", [])
            if hits:
                url = random.choice(hits).get("largeImageURL") or random.choice(hits).get("webformatURL")
                img_r = requests.get(url, timeout=12)
                if img_r.status_code == 200:
                    return BytesIO(img_r.content)
    except Exception as e:
        print(f"Pixabay error: {e}")
    return None


def fetch_image(keyword: str):
    return fetch_unsplash_image(keyword) or fetch_pixabay_image(keyword)


# ─────────────────────────────────────────────────────────────────
#  IMAGE PROCESSING
# ─────────────────────────────────────────────────────────────────
def rounded_image(stream, radius=80) -> BytesIO:
    """Crop to exact ratio, apply rounded mask, return PNG BytesIO."""
    try:
        img = Image.open(stream).convert("RGBA")
        w, h = img.size
        # Crop to 4:3
        target_ratio = 4 / 3
        if w / h > target_ratio:
            new_w = int(h * target_ratio)
            img = img.crop(((w - new_w) // 2, 0, (w + new_w) // 2, h))
        else:
            new_h = int(w / target_ratio)
            img = img.crop((0, (h - new_h) // 2, w, (h + new_h) // 2))

        img = img.resize((800, 600), Image.LANCZOS)
        mask = Image.new("L", (800, 600), 0)
        ImageDraw.Draw(mask).rounded_rectangle([0, 0, 799, 599], radius=radius, fill=255)
        img.putalpha(mask)

        out = BytesIO()
        img.save(out, format="PNG")
        out.seek(0)
        return out
    except Exception as e:
        print(f"rounded_image error: {e}")
        try:
            stream.seek(0)
        except Exception:
            pass
        return stream


def darkened_image(stream, opacity=0.55) -> BytesIO:
    """Full-bleed image darkened for text overlay."""
    try:
        img = Image.open(stream).convert("RGBA")
        img = img.resize((1920, 1080), Image.LANCZOS)
        overlay = Image.new("RGBA", img.size, (0, 0, 0, int(255 * opacity)))
        result = Image.alpha_composite(img, overlay)
        out = BytesIO()
        result.save(out, format="PNG")
        out.seek(0)
        return out
    except Exception as e:
        print(f"darkened_image error: {e}")
        try:
            stream.seek(0)
        except Exception:
            pass
        return stream


# ─────────────────────────────────────────────────────────────────
#  PRIMITIVE HELPERS
# ─────────────────────────────────────────────────────────────────
def bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color: RGBColor, radius: int = 0):
    from pptx.oxml.ns import qn
    from lxml import etree
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius > 0:
        sp = shape._element
        spPr = sp.find(qn("p:spPr"))
        prstGeom = spPr.find(qn("a:prstGeom")) if spPr is not None else None
        if prstGeom is not None:
            spPr.remove(prstGeom)
        pg = etree.SubElement(spPr, qn("a:prstGeom"))
        pg.set("prst", "roundRect")
        avLst = etree.SubElement(pg, qn("a:avLst"))
        gd = etree.SubElement(avLst, qn("a:gd"))
        gd.set("name", "adj")
        gd.set("fmla", f"val {radius}")
    return shape


def txt(slide, text, l, t, w, h, size, color: RGBColor,
        bold=False, italic=False, align=PP_ALIGN.LEFT,
        font="Calibri", wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return tb


def img(slide, stream, l, t, w, h):
    try:
        slide.shapes.add_picture(stream, l, t, w, h)
        return True
    except Exception as e:
        print(f"img error: {e}")
        return False


# ─────────────────────────────────────────────────────────────────
#  SLIDE 0 — COVER  (dark, full bleed image + large title)
# ─────────────────────────────────────────────────────────────────
def build_cover(prs, title, theme, keyword="abstract"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, theme["dark_bg"])
    
    image = fetch_image(keyword)
    if image:
        dark = darkened_image(image, opacity=0.52)
        img(slide, dark, 0, 0, SLIDE_W, SLIDE_H)
    
    # Left thick accent bar
    rect(slide, 0, 0, Inches(0.18), SLIDE_H, theme["accent"])
    
    # Bottom strip
    rect(slide, 0, Inches(6.3), SLIDE_W, Inches(1.2), theme["dark_bg"])
    
    # Title
    txt(slide, title, Inches(0.55), Inches(1.8), Inches(11.5), Inches(3.2), 
        Pt(52), theme["white"], bold=True, font="Calibri", align=PP_ALIGN.LEFT)

# ─────────────────────────────────────────────────────────────────
#  SLIDE 1 — INTRODUCTION  (light bg, description + icon bullets)
# ─────────────────────────────────────────────────────────────────
def build_intro(prs, heading, description, bullets, theme, keyword):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, theme["light_bg"])

    # Left sidebar
    rect(slide, 0, 0, Inches(0.18), SLIDE_H, theme["accent"])

    # Top label chip
    rect(slide, Inches(0.55), Inches(0.38), Inches(1.8), Inches(0.38),
         theme["accent"], radius=20000)
    txt(slide, "INTRODUCTION",
        Inches(0.58), Inches(0.38), Inches(1.8), Inches(0.38),
        Pt(10), theme["white"], bold=True, align=PP_ALIGN.CENTER)

    # Heading
    txt(slide, heading,
        Inches(0.55), Inches(0.92), Inches(7.8), Inches(1.5),
        Pt(38), theme["heading"], bold=True, font="Calibri")

    # Description paragraph
    if description:
        txt(slide, description,
            Inches(0.55), Inches(2.55), Inches(7.6), Inches(1.1),
            Pt(16), theme["body"], font="Calibri")

    # Bullet rows with accent dot
    top = Inches(3.85)
    for bullet in bullets[:4]:
        rect(slide, Inches(0.55), top + Inches(0.13),
             Inches(0.22), Inches(0.22), theme["accent"], radius=50000)
        txt(slide, bullet,
            Inches(0.95), top, Inches(6.9), Inches(0.62),
            Pt(15), theme["body"], font="Calibri")
        top += Inches(0.72)

    # Right image
    image = fetch_image(keyword)
    if image:
        ri = rounded_image(image, radius=60)
        img(slide, ri, Inches(8.6), Inches(0.55), Inches(4.5), Inches(6.4))
    else:
        rect(slide, Inches(8.6), Inches(0.55), Inches(4.5), Inches(6.4),
             theme["card_bg"], radius=30000)


# ─────────────────────────────────────────────────────────────────
#  LAYOUT A — Big stat / hero number left, content right
#  Best for: data-heavy, research, impact slides
# ─────────────────────────────────────────────────────────────────
def build_layout_a(prs, heading, bullets, theme, keyword):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, theme["light_bg"])

    # Top accent bar full width
    rect(slide, 0, 0, SLIDE_W, Inches(0.12), theme["accent"])

    # Full-height left panel (dark)
    rect(slide, 0, 0, Inches(5.1), SLIDE_H, theme["dark_bg"])

    # Image inside left panel
    image = fetch_image(keyword)
    if image:
        ri = rounded_image(image, radius=40)
        img(slide, ri, Inches(0.25), Inches(0.5), Inches(4.6), Inches(4.5))
    else:
        rect(slide, Inches(0.25), Inches(0.5), Inches(4.6), Inches(4.5),
             theme["card_bg"], radius=20000)

    # Heading label on left panel bottom
    txt(slide, heading,
        Inches(0.3), Inches(5.2), Inches(4.6), Inches(1.9),
        Pt(22), theme["accent_soft"], bold=True, font="Calibri")

    # Right: bullets with numbered chips
    txt(slide, "Key Points",
        Inches(5.4), Inches(0.35), Inches(7.6), Inches(0.6),
        Pt(12), theme["muted"], bold=True, font="Calibri")

    top = Inches(1.05)
    for i, bullet in enumerate(bullets[:5]):
        # Number chip
        rect(slide, Inches(5.4), top, Inches(0.42), Inches(0.42),
             theme["accent"], radius=50000)
        txt(slide, str(i + 1),
            Inches(5.4), top, Inches(0.42), Inches(0.42),
            Pt(13), theme["white"], bold=True, align=PP_ALIGN.CENTER)
        # Bullet text
        txt(slide, bullet,
            Inches(6.0), top, Inches(7.0), Inches(0.6),
            Pt(16), theme["body"], font="Calibri")
        top += Inches(1.08)


# ─────────────────────────────────────────────────────────────────
#  LAYOUT B — Full bleed image, content overlay bottom half
#  Best for: dramatic, impactful statements
# ─────────────────────────────────────────────────────────────────
def build_layout_b(prs, heading, bullets, theme, keyword):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, theme["dark_bg"])

    image = fetch_image(keyword)
    if image:
        dark = darkened_image(image, opacity=0.45)
        img(slide, dark, 0, 0, SLIDE_W, SLIDE_H)

    # Dark overlay bottom 60%
    rect(slide, 0, Inches(2.9), SLIDE_W, Inches(4.6), theme["dark_bg"])

    # Accent top bar
    rect(slide, 0, 0, SLIDE_W, Inches(0.1), theme["accent"])

    # Heading large, white — sits right on the overlay start
    txt(slide, heading,
        Inches(0.7), Inches(3.0), Inches(11.5), Inches(1.5),
        Pt(40), theme["white"], bold=True, font="Calibri")

    # Two-column bullets
    mid = len(bullets[:4]) // 2 or 2
    col1 = bullets[:mid]
    col2 = bullets[mid:4]

    top = Inches(4.65)
    for b in col1:
        txt(slide, f"→  {b}",
            Inches(0.7), top, Inches(5.7), Inches(0.65),
            Pt(15), theme["accent_soft"], font="Calibri")
        top += Inches(0.75)

    top = Inches(4.65)
    for b in col2:
        txt(slide, f"→  {b}",
            Inches(6.9), top, Inches(5.7), Inches(0.65),
            Pt(15), theme["accent_soft"], font="Calibri")
        top += Inches(0.75)


# ─────────────────────────────────────────────────────────────────
#  LAYOUT C — 2x2 card grid  (no image — content is the visual)
#  Best for: 4 distinct concepts, comparisons
# ─────────────────────────────────────────────────────────────────
def build_layout_c(prs, heading, bullets, theme, keyword):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, theme["light_bg"])

    rect(slide, 0, 0, SLIDE_W, Inches(0.12), theme["accent"])

    txt(slide, heading,
        Inches(0.6), Inches(0.28), Inches(11.5), Inches(0.9),
        Pt(36), theme["heading"], bold=True, font="Calibri")

    # 4 cards in 2x2 grid
    card_data = bullets[:4]
    positions = [
        (Inches(0.5),  Inches(1.45)),
        (Inches(6.95), Inches(1.45)),
        (Inches(0.5),  Inches(4.3)),
        (Inches(6.95), Inches(4.3)),
    ]
    card_w, card_h = Inches(6.2), Inches(2.6)

    accent_colors = [theme["accent"], theme["dark_bg"],
                     theme["dark_bg"], theme["accent"]]

    for i, (btext, (cl, ct)) in enumerate(zip(card_data, positions)):
        c = accent_colors[i % len(accent_colors)]
        rect(slide, cl, ct, card_w, card_h, c, radius=15000)
        # Number
        txt(slide, f"0{i+1}",
            cl + Inches(0.25), ct + Inches(0.2), Inches(0.8), Inches(0.55),
            Pt(22), theme["accent_soft"] if c == theme["dark_bg"] else theme["white"],
            bold=True, font="Calibri")
        # Text
        txt(slide, btext,
            cl + Inches(0.25), ct + Inches(0.9), card_w - Inches(0.5), Inches(1.5),
            Pt(15), theme["white"], font="Calibri")


# ─────────────────────────────────────────────────────────────────
#  LAYOUT D — Timeline / process steps
#  Best for: sequential info, how-something-works
# ─────────────────────────────────────────────────────────────────
def build_layout_d(prs, heading, bullets, theme, keyword):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, theme["light_bg"])

    rect(slide, 0, 0, SLIDE_W, Inches(0.12), theme["accent"])

    txt(slide, heading,
        Inches(0.6), Inches(0.28), Inches(10.0), Inches(0.9),
        Pt(36), theme["heading"], bold=True, font="Calibri")

    # Horizontal timeline connector line
    n = min(len(bullets), 5)
    step_w = Inches(13.33 / (n + 0.5))
    start_x = Inches(0.5)
    line_y = Inches(2.6)

    rect(slide, start_x, line_y + Inches(0.15),
         SLIDE_W - Inches(1.0), Inches(0.05), theme["accent_soft"])

    for i, bullet in enumerate(bullets[:5]):
        cx = start_x + i * step_w + step_w * 0.1

        # Circle node on timeline
        rect(slide, cx, line_y, Inches(0.36), Inches(0.36),
             theme["accent"], radius=50000)
        txt(slide, str(i + 1),
            cx, line_y, Inches(0.36), Inches(0.36),
            Pt(12), theme["white"], bold=True, align=PP_ALIGN.CENTER)

        # Step label above
        txt(slide, f"STEP {i+1:02d}",
            cx - Inches(0.3), line_y - Inches(0.55), Inches(1.4), Inches(0.4),
            Pt(10), theme["muted"], bold=True, font="Calibri")

        # Text below — taller box to fill space
        txt(slide, bullet,
            cx - Inches(0.3), line_y + Inches(0.55), step_w - Inches(0.1), Inches(3.8),
            Pt(14), theme["body"], font="Calibri", wrap=True)

    # Small image top-right
    image = fetch_image(keyword)
    if image:
        ri = rounded_image(image, radius=30)
        img(slide, ri, Inches(10.5), Inches(0.28), Inches(2.6), Inches(1.9))


# ─────────────────────────────────────────────────────────────────
#  LAYOUT E — Split: image right, icon-list left
#  Best for: general content, clean readable slides
# ─────────────────────────────────────────────────────────────────
def build_layout_e(prs, heading, bullets, theme, keyword):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, theme["light_bg"])

    rect(slide, 0, 0, SLIDE_W, Inches(0.12), theme["accent"])

    # Heading
    txt(slide, heading,
        Inches(0.55), Inches(0.28), Inches(7.8), Inches(1.2),
        Pt(36), theme["heading"], bold=True, font="Calibri")

    # Icon bullets — square accent chip + text
    top = Inches(1.65)
    for bullet in bullets[:5]:
        rect(slide, Inches(0.55), top + Inches(0.06),
             Inches(0.28), Inches(0.28), theme["accent"], radius=8000)
        txt(slide, bullet,
            Inches(1.05), top, Inches(6.9), Inches(0.6),
            Pt(16), theme["body"], font="Calibri")
        top += Inches(0.9)

    # Right image tall
    image = fetch_image(keyword)
    if image:
        ri = rounded_image(image, radius=50)
        img(slide, ri, Inches(8.55), Inches(0.3), Inches(4.55), Inches(6.9))
    else:
        rect(slide, Inches(8.55), Inches(0.3), Inches(4.55), Inches(6.9),
             theme["card_bg"], radius=20000)


# ─────────────────────────────────────────────────────────────────
#  LAYOUT F — Dark left panel + light right content
#  Best for: contrast, premium feel, emphasis slides
# ─────────────────────────────────────────────────────────────────
def build_layout_f(prs, heading, bullets, theme, keyword):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, theme["light_bg"])

    # Full-height dark left panel
    rect(slide, 0, 0, Inches(4.8), SLIDE_H, theme["dark_bg"])

    # Accent bottom strip on left panel
    rect(slide, 0, Inches(6.8), Inches(4.8), Inches(0.7), theme["accent"])

    # Image inside left panel
    image = fetch_image(keyword)
    if image:
        ri = rounded_image(image, radius=30)
        img(slide, ri, Inches(0.25), Inches(0.4), Inches(4.3), Inches(4.8))
    else:
        rect(slide, Inches(0.25), Inches(0.4), Inches(4.3), Inches(4.8),
             theme["card_bg"], radius=15000)

    # Heading in left panel bottom area
    txt(slide, heading,
        Inches(0.3), Inches(5.35), Inches(4.3), Inches(1.3),
        Pt(20), theme["accent_soft"], bold=True, font="Calibri")

    # Right: large heading label + bullets
    txt(slide, "Insights",
        Inches(5.2), Inches(0.35), Inches(7.8), Inches(0.55),
        Pt(12), theme["muted"], bold=True, font="Calibri")

    top = Inches(1.1)
    for bullet in bullets[:5]:
        rect(slide, Inches(5.2), top + Inches(0.15),
             Inches(0.08), Inches(0.3), theme["accent"])
        txt(slide, bullet,
            Inches(5.55), top, Inches(7.5), Inches(0.65),
            Pt(16), theme["body"], font="Calibri")
        top += Inches(1.0)


# ─────────────────────────────────────────────────────────────────
#  CONCLUSION SLIDE — dark, summary-focused
# ─────────────────────────────────────────────────────────────────
def build_conclusion(prs, heading, description, bullets, theme, keyword):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, theme["dark_bg"])

    # Background image faint
    image = fetch_image(keyword)
    if image:
        dark = darkened_image(image, opacity=0.72)
        img(slide, dark, 0, 0, SLIDE_W, SLIDE_H)

    # Left accent bar
    rect(slide, 0, 0, Inches(0.18), SLIDE_H, theme["accent"])

    # "Conclusion" label chip
    rect(slide, Inches(0.55), Inches(0.38), Inches(1.9), Inches(0.38),
         theme["accent"], radius=20000)
    txt(slide, "CONCLUSION",
        Inches(0.58), Inches(0.38), Inches(1.9), Inches(0.38),
        Pt(10), theme["white"], bold=True, align=PP_ALIGN.CENTER)

    # Heading
    txt(slide, heading,
        Inches(0.55), Inches(0.92), Inches(11.5), Inches(1.3),
        Pt(38), theme["white"], bold=True, font="Calibri")

    # Description
    if description:
        txt(slide, description,
            Inches(0.55), Inches(2.35), Inches(11.5), Inches(0.9),
            Pt(16), theme["accent_soft"], font="Calibri", italic=True)

    # Key takeaway bullets - 2 columns
    col1 = bullets[:2]
    col2 = bullets[2:4]

    top = Inches(3.45)
    for b in col1:
        rect(slide, Inches(0.55), top + Inches(0.1),
             Inches(0.22), Inches(0.22), theme["accent"], radius=50000)
        txt(slide, b, Inches(0.95), top, Inches(5.5), Inches(0.65),
            Pt(15), theme["white"], font="Calibri")
        top += Inches(0.82)

    top = Inches(3.45)
    for b in col2:
        rect(slide, Inches(6.8), top + Inches(0.1),
             Inches(0.22), Inches(0.22), theme["accent"], radius=50000)
        txt(slide, b, Inches(7.2), top, Inches(5.5), Inches(0.65),
            Pt(15), theme["white"], font="Calibri")
        top += Inches(0.82)

    # SlideBot credit



# ─────────────────────────────────────────────────────────────────
#  THANK YOU SLIDE
# ─────────────────────────────────────────────────────────────────
def build_thankyou(prs, theme, is_premium: bool = False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, theme["dark_bg"])
    
    # Left accent bar
    rect(slide, 0, 0, Inches(0.18), SLIDE_H, theme["accent"])
    
    # Main "Thank You" text
    txt(slide, "Thank You", Inches(0.6), Inches(2.2), Inches(11.5), Inches(2.0), 
        Pt(72), theme["white"], bold=True, font="Calibri", align=PP_ALIGN.LEFT)
    
    # Watermark: Only for FREE users, very subtle at bottom right
    if not is_premium:
        # Small, low-opacity text at bottom right (looks like a soft sticker)
        txt(slide, "Created with SlideBot", 
            Inches(9.8), Inches(6.9), Inches(3.2), Inches(0.4), 
            Pt(10), theme["muted"], italic=True, font="Calibri", 
            align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────
#  MAIN BUILDER
# ─────────────────────────────────────────────────────────────────
CONTENT_LAYOUTS = [
    build_layout_a,
    build_layout_b,
    build_layout_c,
    build_layout_d,
    build_layout_e,
    build_layout_f,
]


def build_presentation(slide_data: dict, theme_name: str = "classic",
                        is_premium: bool = False) -> str:
    theme  = THEMES.get(theme_name, THEMES["classic"])
    slides = slide_data.get("slides", [])
    title  = slide_data.get("title", "My Presentation")

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    if not slides:
        slides = []

    # ── Cover
    cover_kw = slides[0].get("image_keyword", "abstract") if slides else "abstract"
    build_cover(prs, title, theme, cover_kw)

    # ── Intro (first slide in data)
    if slides:
        s0 = slides[0]
        build_intro(
            prs,
            s0.get("heading", "Introduction"),
            s0.get("description", ""),
            s0.get("bullets", []),
            theme,
            s0.get("image_keyword", "teamwork"),
        )

    # ── Content slides (everything except first and last)
    content = slides[1:-1] if len(slides) > 2 else []
    for idx, s in enumerate(content):
        fn = CONTENT_LAYOUTS[idx % len(CONTENT_LAYOUTS)]
        fn(prs, s.get("heading", ""), s.get("bullets", []),
           theme, s.get("image_keyword", "business"))

    # ── Conclusion (last slide in data)
    if len(slides) > 1:
        sl = slides[-1]
        build_conclusion(
            prs,
            sl.get("heading", "Conclusion & Key Takeaways"),
            sl.get("description", ""),
            sl.get("bullets", []),
            theme,
            sl.get("image_keyword", "success"),
        )

    # ── Thank You
    build_thankyou(prs, theme, is_premium=is_premium)

    filename = f"slidebot_{uuid.uuid4().hex[:8]}.pptx"
    filepath = os.path.join("outputs", filename)
    os.makedirs("outputs", exist_ok=True)
    prs.save(filepath)
    print(f"✅ Saved: {filepath}")
    return filepath
