import os
import uuid
import random
import requests
from io import BytesIO

from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from PIL import Image, ImageDraw

load_dotenv()

UNSPLASH_KEY = os.getenv("UNSPLASH_ACCESS_KEY")
PIXABAY_KEY = os.getenv("PIXABAY_API_KEY")

# ─────────────────────────────────────────────────────────────────
# COLOR SYSTEM
# ─────────────────────────────────────────────────────────────────
MAGAZINE_COLORS = {
    "navy": RGBColor(0x1E, 0x27, 0x61),
    "red": RGBColor(0xC0, 0x39, 0x2B),
    "green": RGBColor(0x2C, 0x5F, 0x2D),
    "purple": RGBColor(0x6B, 0x2D, 0x8B),
    "orange": RGBColor(0xFF, 0x6B, 0x35),
    "gold": RGBColor(0xB7, 0x79, 0x0F),
    "teal": RGBColor(0x02, 0x80, 0x90),
    "pink": RGBColor(0xE9, 0x4C, 0x7D),
}

COLOR_MAP = {
    "navy": "navy", "blue": "navy", "indigo": "navy",
    "red": "red", "crimson": "red", "scarlet": "red",
    "green": "green", "forest": "green", "emerald": "green",
    "purple": "purple", "violet": "purple", "plum": "purple",
    "orange": "orange", "amber": "orange",
    "gold": "gold", "yellow": "gold", "mustard": "gold",
    "teal": "teal", "cyan": "teal", "aqua": "teal",
    "pink": "pink", "rose": "pink", "magenta": "pink",
}

MAG = {
    "black": RGBColor(0x0D, 0x0D, 0x0D),
    "off_white": RGBColor(0xF9, 0xF7, 0xF4),
    "light_gray": RGBColor(0xE8, 0xE6, 0xE3),
    "mid_gray": RGBColor(0x88, 0x88, 0x88),
    "dark_gray": RGBColor(0x2D, 0x2D, 0x2D),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
}

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ─────────────────────────────────────────────────────────────────
# IMAGE HELPERS
# ─────────────────────────────────────────────────────────────────
def fetch_image(keyword):
    try:
        if UNSPLASH_KEY:
            r = requests.get(
                "https://api.unsplash.com/photos/random",
                params={"query": keyword, "orientation": "landscape",
                        "content_filter": "high", "client_id": UNSPLASH_KEY},
                timeout=12)
            if r.status_code == 200:
                url = r.json()["urls"]["regular"]
                ir = requests.get(url, timeout=12)
                if ir.status_code == 200:
                    return BytesIO(ir.content)
    except Exception as e:
        print(f"Unsplash: {e}")

    try:
        if PIXABAY_KEY:
            r = requests.get(
                "https://pixabay.com/api/",
                params={"key": PIXABAY_KEY, "q": keyword, "image_type": "photo",
                        "orientation": "horizontal", "safesearch": "true", "per_page": 5},
                timeout=12)
            if r.status_code == 200:
                hits = r.json().get("hits", [])
                if hits:
                    url = random.choice(hits).get("largeImageURL")
                    ir = requests.get(url, timeout=12)
                    if ir.status_code == 200:
                        return BytesIO(ir.content)
    except Exception as e:
        print(f"Pixabay: {e}")
    return None


def process_image(stream, w_px=1200, h_px=900, radius=0):
    try:
        img = Image.open(stream).convert("RGBA")
        iw, ih = img.size
        target = w_px / h_px
        if iw / ih > target:
            new_w = int(ih * target)
            img = img.crop(((iw-new_w)//2, 0, (iw+new_w)//2, ih))
        else:
            new_h = int(iw / target)
            img = img.crop((0, (ih-new_h)//2, iw, (ih+new_h)//2))
        img = img.resize((w_px, h_px), Image.LANCZOS)
        if radius:
            mask = Image.new("L", (w_px, h_px), 0)
            ImageDraw.Draw(mask).rounded_rectangle(
                [0, 0, w_px-1, h_px-1], radius=radius, fill=255)
            img.putalpha(mask)
        out = BytesIO()
        img.save(out, format="PNG")
        out.seek(0)
        return out
    except Exception as e:
        print(f"process_image: {e}")
        try: stream.seek(0)
        except: pass
        return stream


def darkened(stream, opacity=0.5):
    try:
        img = Image.open(stream).convert("RGBA")
        img = img.resize((1920, 1080), Image.LANCZOS)
        overlay = Image.new("RGBA", img.size, (0, 0, 0, int(255*opacity)))
        out = BytesIO()
        Image.alpha_composite(img, overlay).save(out, format="PNG")
        out.seek(0)
        return out
    except Exception as e:
        print(f"darkened: {e}")
        try: stream.seek(0)
        except: pass
        return stream


# ─────────────────────────────────────────────────────────────────
# DRAWING PRIMITIVES
# ─────────────────────────────────────────────────────────────────
def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, l, t, w, h, color, radius=0):
    from pptx.oxml.ns import qn
    from lxml import etree
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius:
        sp = shape._element
        spPr = sp.find(qn("p:spPr"))
        pg = spPr.find(qn("a:prstGeom"))
        if pg is not None:
            spPr.remove(pg)
        pg = etree.SubElement(spPr, qn("a:prstGeom"))
        pg.set("prst", "roundRect")
        av = etree.SubElement(pg, qn("a:avLst"))
        gd = etree.SubElement(av, qn("a:gd"))
        gd.set("name", "adj")
        gd.set("fmla", f"val {radius}")
    return shape


def txt(slide, text, l, t, w, h, size, color,
        bold=False, italic=False, align=PP_ALIGN.LEFT, font="Georgia"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return tb


def pic(slide, stream, l, t, w, h):
    try:
        slide.shapes.add_picture(stream, l, t, w, h)
        return True
    except Exception as e:
        print(f"pic: {e}")
        return False


# ─────────────────────────────────────────────────────────────────
# COVER — No "Generated by SlideBot"
# ─────────────────────────────────────────────────────────────────
def mag_cover(prs, title, accent, keyword):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, MAG["black"])

    image = fetch_image(keyword)
    if image:
        dark = darkened(image, opacity=0.6)
        pic(slide, dark, 0, 0, SLIDE_W, SLIDE_H)

    # Thick left accent strip
    box(slide, 0, 0, Inches(0.55), SLIDE_H, accent)

    # Top label bar
    box(slide, Inches(0.55), 0, SLIDE_W, Inches(0.65), MAG["black"])
    txt(slide, "S L I D E B O T P R E S E N T A T I O N",
        Inches(0.85), Inches(0.1), Inches(11.0), Inches(0.48),
        Pt(11), MAG["mid_gray"], font="Calibri")

    # Giant editorial title
    txt(slide, title,
        Inches(0.85), Inches(1.0), Inches(10.5), Inches(4.6),
        Pt(56), MAG["white"], bold=True, font="Georgia")

    # Accent rule
    box(slide, Inches(0.85), Inches(5.8), Inches(3.5), Inches(0.07), accent)


# ─────────────────────────────────────────────────────────────────
# INTRO (unchanged)
# ─────────────────────────────────────────────────────────────────
def mag_intro(prs, heading, description, bullets, accent, keyword):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, MAG["off_white"])

    # Left black column with image
    box(slide, 0, 0, Inches(4.5), SLIDE_H, MAG["black"])
    image = fetch_image(keyword)
    if image:
        proc = process_image(image, 800, 1080)
        pic(slide, proc, 0, 0, Inches(4.5), SLIDE_H)
    box(slide, 0, Inches(4.8), Inches(4.5), Inches(2.7), MAG["black"])

    txt(slide, "INTRODUCTION",
        Inches(0.2), Inches(5.15), Inches(4.0), Inches(0.45),
        Pt(10), accent, bold=True, font="Calibri")

    txt(slide, heading,
        Inches(0.25), Inches(5.7), Inches(4.1), Inches(1.55),
        Pt(21), MAG["white"], bold=True, font="Georgia")

    # Right: pull quote
    if description:
        txt(slide, "\u201C",
            Inches(4.8), Inches(0.2), Inches(1.5), Inches(1.4),
            Pt(96), accent, bold=True, font="Georgia")
        txt(slide, description,
            Inches(4.8), Inches(1.25), Inches(8.2), Inches(1.85),
            Pt(19), MAG["dark_gray"], italic=True, font="Georgia")

    box(slide, Inches(4.8), Inches(3.3), Inches(8.2), Inches(0.03), MAG["light_gray"])

    top = Inches(3.55)
    for bullet in bullets[:4]:
        box(slide, Inches(4.8), top + Inches(0.12),
            Inches(0.18), Inches(0.18), accent, radius=50000)
        txt(slide, bullet,
            Inches(5.2), top, Inches(7.8), Inches(0.6),
            Pt(15), MAG["dark_gray"], font="Calibri")
        top += Inches(0.72)


# ─────────────────────────────────────────────────────────────────
# LAYOUTS 1 to 5 (unchanged - keeping them as they were)
# ─────────────────────────────────────────────────────────────────
def mag_layout_1(prs, heading, bullets, accent, keyword):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, MAG["off_white"])

    image = fetch_image(keyword)
    if image:
        proc = process_image(image, 1920, 580)
        pic(slide, proc, 0, 0, SLIDE_W, Inches(4.0))

    # Heading overlay strip
    box(slide, 0, Inches(2.75), SLIDE_W, Inches(1.25), MAG["black"])
    txt(slide, heading,
        Inches(0.6), Inches(2.82), Inches(12.0), Inches(1.1),
        Pt(34), MAG["white"], bold=True, font="Georgia")

    # Tag top-left
    box(slide, Inches(0.5), Inches(0.28), Inches(0.06), Inches(0.55), accent)
    txt(slide, "KEY INSIGHT",
        Inches(0.7), Inches(0.3), Inches(3.0), Inches(0.4),
        Pt(10), accent, bold=True, font="Calibri")

    col1 = bullets[:2]
    col2 = bullets[2:4]
    top = Inches(4.28)
    for b in col1:
        txt(slide, f"— {b}", Inches(0.6), top, Inches(5.8), Inches(0.75),
            Pt(15), MAG["dark_gray"], font="Calibri")
        top += Inches(0.85)

    top = Inches(4.28)
    for b in col2:
        txt(slide, f"— {b}", Inches(7.0), top, Inches(5.8), Inches(0.75),
            Pt(15), MAG["dark_gray"], font="Calibri")
        top += Inches(0.85)

    box(slide, 0, Inches(7.35), SLIDE_W, Inches(0.15), accent)


def mag_layout_2(prs, heading, bullets, accent, keyword):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, MAG["off_white"])

    box(slide, 0, 0, SLIDE_W, Inches(0.08), accent)

    txt(slide, heading,
        Inches(0.6), Inches(0.22), Inches(11.5), Inches(1.0),
        Pt(34), MAG["black"], bold=True, font="Georgia")

    box(slide, Inches(0.6), Inches(1.38), Inches(12.0), Inches(0.03), MAG["light_gray"])

    # Big number block
    box(slide, Inches(0.5), Inches(1.58), Inches(3.8), Inches(5.7), accent)
    txt(slide, str(len(bullets)),
        Inches(0.5), Inches(1.95), Inches(3.8), Inches(2.5),
        Pt(120), MAG["white"], bold=True, font="Georgia", align=PP_ALIGN.CENTER)
    txt(slide, "KEY\nPOINTS",
        Inches(0.5), Inches(4.55), Inches(3.8), Inches(1.2),
        Pt(18), MAG["white"], bold=True, font="Calibri", align=PP_ALIGN.CENTER)

    # Editorial list right
    top = Inches(1.62)
    for i, bullet in enumerate(bullets[:5]):
        box(slide, Inches(4.7), top, Inches(8.3), Inches(0.025), MAG["light_gray"])
        txt(slide, f"{i+1:02d}",
            Inches(4.7), top + Inches(0.08), Inches(0.75), Inches(0.78),
            Pt(22), accent, bold=True, font="Georgia")
        txt(slide, bullet,
            Inches(5.65), top + Inches(0.1), Inches(7.3), Inches(0.78),
            Pt(15), MAG["dark_gray"], font="Calibri")
        top += Inches(1.08)


def mag_layout_3(prs, heading, bullets, accent, keyword):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, MAG["off_white"])

    image = fetch_image(keyword)
    if image:
        proc = process_image(image, 900, 1080)
        pic(slide, proc, Inches(6.8), 0, Inches(6.53), SLIDE_H)

    box(slide, Inches(6.75), Inches(0.4), Inches(0.04), Inches(6.7), MAG["light_gray"])

    txt(slide, "— ANALYSIS",
        Inches(0.5), Inches(0.28), Inches(5.0), Inches(0.45),
        Pt(11), accent, bold=True, font="Calibri")

    txt(slide, heading,
        Inches(0.5), Inches(0.82), Inches(6.0), Inches(2.4),
        Pt(38), MAG["black"], bold=True, font="Georgia")

    box(slide, Inches(0.5), Inches(3.38), Inches(1.8), Inches(0.07), accent)

    top = Inches(3.62)
    for bullet in bullets[:4]:
        txt(slide, bullet,
            Inches(0.5), top, Inches(6.0), Inches(0.7),
            Pt(15), MAG["dark_gray"], font="Calibri")
        box(slide, Inches(0.5), top + Inches(0.73),
            Inches(6.0), Inches(0.02), MAG["light_gray"])
        top += Inches(0.9)


def mag_layout_4(prs, heading, bullets, accent, keyword):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, MAG["black"])

    image = fetch_image(keyword)
    if image:
        dark = darkened(image, opacity=0.65)
        pic(slide, dark, 0, 0, SLIDE_W, SLIDE_H)

    box(slide, 0, 0, SLIDE_W, Inches(0.08), accent)
    box(slide, 0, Inches(7.42), SLIDE_W, Inches(0.08), accent)

    txt(slide, heading,
        Inches(1.0), Inches(1.1), Inches(11.33), Inches(2.6),
        Pt(46), MAG["white"], bold=True, font="Georgia", align=PP_ALIGN.CENTER)

    box(slide, Inches(5.5), Inches(3.88), Inches(2.33), Inches(0.07), accent)

    col1 = bullets[:2]
    col2 = bullets[2:4]
    top = Inches(4.18)
    for b in col1:
        txt(slide, f"\u25C6 {b}", Inches(0.8), top, Inches(5.6), Inches(0.72),
            Pt(15), MAG["white"], font="Calibri")
        top += Inches(0.85)

    top = Inches(4.18)
    for b in col2:
        txt(slide, f"\u25C6 {b}", Inches(7.0), top, Inches(5.6), Inches(0.72),
            Pt(15), MAG["white"], font="Calibri")
        top += Inches(0.85)


def mag_layout_5(prs, heading, bullets, accent, keyword):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, MAG["off_white"])

    box(slide, 0, 0, SLIDE_W, Inches(0.08), accent)

    txt(slide, heading,
        Inches(0.6), Inches(0.2), Inches(11.5), Inches(0.95),
        Pt(34), MAG["black"], bold=True, font="Georgia")

    box(slide, Inches(0.6), Inches(1.25), Inches(12.0), Inches(0.03), MAG["light_gray"])

    cols = [bullets[:2], bullets[2:4], bullets[4:5]]
    labels = ["01", "02", "03"]
    col_x = [Inches(0.5), Inches(4.72), Inches(8.94)]
    col_w = Inches(3.9)
    image = fetch_image(keyword)

    for i, (cx, label, col_bullets) in enumerate(zip(col_x, labels, cols)):
        txt(slide, label, cx, Inches(1.38), col_w, Inches(1.05),
            Pt(52), accent, bold=True, font="Georgia")

        if i == 1 and image:
            proc = process_image(image, 600, 400, radius=20)
            pic(slide, proc, cx, Inches(2.52), col_w, Inches(2.18))
            top = Inches(4.85)
        else:
            top = Inches(2.55)

        for b in col_bullets:
            txt(slide, b, cx, top, col_w - Inches(0.12), Inches(1.05),
                Pt(14), MAG["dark_gray"], font="Calibri")
            top += Inches(1.12)

        if i < 2:
            box(slide, cx + col_w + Inches(0.06), Inches(1.38),
                Inches(0.025), Inches(5.85), MAG["light_gray"])

    box(slide, 0, Inches(7.35), SLIDE_W, Inches(0.15), accent)


# ─────────────────────────────────────────────────────────────────
# CONCLUSION — No "Created with SlideBot"
# ─────────────────────────────────────────────────────────────────
def mag_conclusion(prs, heading, description, bullets, accent, keyword):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, MAG["black"])

    image = fetch_image(keyword)
    if image:
        dark = darkened(image, opacity=0.7)
        pic(slide, dark, 0, 0, SLIDE_W, SLIDE_H)

    box(slide, 0, 0, Inches(0.55), SLIDE_H, accent)

    box(slide, Inches(0.8), Inches(0.32), Inches(2.1), Inches(0.44), accent)
    txt(slide, "CONCLUSION",
        Inches(0.83), Inches(0.32), Inches(2.1), Inches(0.44),
        Pt(11), MAG["white"], bold=True, font="Calibri", align=PP_ALIGN.CENTER)

    txt(slide, heading,
        Inches(0.8), Inches(0.9), Inches(11.5), Inches(1.65),
        Pt(42), MAG["white"], bold=True, font="Georgia")

    if description:
        txt(slide, f"\u201C{description}\u201D",
            Inches(0.8), Inches(2.68), Inches(11.5), Inches(1.05),
            Pt(17), RGBColor(0xCC, 0xCC, 0xCC), italic=True, font="Georgia")

    box(slide, Inches(0.8), Inches(3.88), Inches(11.5), Inches(0.04), accent)

    col1 = bullets[:2]
    col2 = bullets[2:4]
    top = Inches(4.1)
    for b in col1:
        box(slide, Inches(0.8), top + Inches(0.12),
            Inches(0.2), Inches(0.2), accent, radius=50000)
        txt(slide, b, Inches(1.2), top, Inches(5.5), Inches(0.65),
            Pt(15), MAG["white"], font="Calibri")
        top += Inches(0.85)

    top = Inches(4.1)
    for b in col2:
        box(slide, Inches(7.1), top + Inches(0.12),
            Inches(0.2), Inches(0.2), accent, radius=50000)
        txt(slide, b, Inches(7.5), top, Inches(5.5), Inches(0.65),
            Pt(15), MAG["white"], font="Calibri")
        top += Inches(0.85)


# ─────────────────────────────────────────────────────────────────
# THANK YOU — Subtle for Free, Clean for Premium
# ─────────────────────────────────────────────────────────────────
def mag_thankyou(prs, accent, is_premium: bool = False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, MAG["black"])
    box(slide, 0, 0, Inches(0.55), SLIDE_H, accent)
    
    txt(slide, "Thank\nYou",
        Inches(0.9), Inches(1.5), Inches(10.0), Inches(4.2),
        Pt(86), MAG["white"], bold=True, font="Georgia")
    
    box(slide, Inches(0.9), Inches(5.75), Inches(4.0), Inches(0.07), accent)
    
    # Only show for free users - small and subtle at bottom right
    if not is_premium:
        txt(slide, "Created with SlideBot",
            Inches(9.2), Inches(6.85), Inches(3.8), Inches(0.35),
            Pt(9.5), MAG["mid_gray"], italic=True, font="Calibri",
            align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────────────────────────
# COLOR RESOLVER
# ─────────────────────────────────────────────────────────────────
def resolve_accent(color_input: str) -> RGBColor:
    if not color_input:
        return MAGAZINE_COLORS["navy"]
    c = color_input.strip().lower().lstrip("#")
    if len(c) == 6:
        try:
            return RGBColor(int(c[0:2],16), int(c[2:4],16), int(c[4:6],16))
        except ValueError:
            pass
    key = COLOR_MAP.get(c)
    if key:
        return MAGAZINE_COLORS[key]
    return MAGAZINE_COLORS["navy"]


# ─────────────────────────────────────────────────────────────────
# MAIN BUILDER
# ─────────────────────────────────────────────────────────────────
CONTENT_LAYOUTS = [
    mag_layout_1, mag_layout_2, mag_layout_3,
    mag_layout_4, mag_layout_5,
]


def build_magazine(slide_data: dict, color_input: str = "navy", is_premium: bool = False) -> str:
    accent = resolve_accent(color_input)
    slides = slide_data.get("slides", [])
    title = slide_data.get("title", "My Presentation")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    cover_kw = slides[0].get("image_keyword", "abstract") if slides else "abstract"
    mag_cover(prs, title, accent, cover_kw)

    if slides:
        s0 = slides[0]
        mag_intro(prs, s0.get("heading", "Introduction"),
                  s0.get("description", ""),
                  s0.get("bullets", []),
                  accent, s0.get("image_keyword", "people"))

    content = slides[1:-1] if len(slides) > 2 else []
    for idx, s in enumerate(content):
        fn = CONTENT_LAYOUTS[idx % len(CONTENT_LAYOUTS)]
        fn(prs, s.get("heading", ""), s.get("bullets", []),
           accent, s.get("image_keyword", "business"))

    if len(slides) > 1:
        sl = slides[-1]
        mag_conclusion(prs, sl.get("heading", "Conclusion"),
                       sl.get("description", ""),
                       sl.get("bullets", []),
                       accent, sl.get("image_keyword", "success"))

    # Thank You with premium flag
    mag_thankyou(prs, accent, is_premium=is_premium)

    filename = f"slidebot_mag_{uuid.uuid4().hex[:8]}.pptx"
    filepath = os.path.join("outputs", filename)
    os.makedirs("outputs", exist_ok=True)
    prs.save(filepath)
    print(f"✅ Magazine pack saved: {filepath}")
    return filepath
